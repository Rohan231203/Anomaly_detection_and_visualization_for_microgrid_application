from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_slide(prs, title, content, bg_color=(255, 255, 255)):
    slide_layout = prs.slide_layouts[5]  # Use a blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)

    # Add title shape
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_paragraph = title_frame.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(36)  # Slightly smaller font for title
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White color
    title_paragraph.alignment = PP_ALIGN.CENTER

    # Add content shape
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(4.5))  # Move closer to bottom
    content_frame = content_box.text_frame
    content_frame.word_wrap = True  # Enable word wrap

    # Add bullet points for the content
    for line in content.split('\n'):
        content_paragraph = content_frame.add_paragraph()
        content_paragraph.text = line
        content_paragraph.font.size = Pt(20)  # Slightly smaller font for content
        content_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color
        content_paragraph.space_after = Pt(14)  # Space after each paragraph

    # Add a decorative shape (footer style)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue color
    shape.line.color.rgb = RGBColor(0, 51, 102)

    # Add a decorative border
    border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(9), Inches(5.5))
    border_shape.line.color.rgb = RGBColor(0, 51, 102)  # Border color
    border_shape.line.width = Pt(2)  # Border width


# Create a presentation object
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]  # Use the title layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Comparison of Electricity Tariffs"
subtitle.text = "A Study of Electricity Tariffs in 5 Countries Including Uttar Pradesh, India"

# Country slides with specific designs and colors
countries = [
    {"name": "Japan", "tariff": "XX.XX USD per kWh"},
    {"name": "USA", "tariff": "XX.XX USD per kWh"},
    {"name": "Canada", "tariff": "XX.XX CAD per kWh"},
    {"name": "UK", "tariff": "XX.XX GBP per kWh"},
    {"name": "Uttar Pradesh, India", "tariff": "XX.XX INR per kWh"},
]

bg_colors = [(173, 216, 230), (255, 182, 193), (221, 160, 221), (240, 230, 140),
             (224, 255, 255)]  # Light colors for slides

for i, country in enumerate(countries):
    create_slide(prs, f"Electricity Tariff in {country['name']}",
                 f"Tariff: {country['tariff']}\n\nExample Calculation:\n"
                 f"For 100 kWh of power consumption:\n"
                 f"Bill = {country['tariff']} * 100", bg_color=bg_colors[i % len(bg_colors)])

# Save the presentation
prs.save("Electricity_Tariffs_Comparison.pptx")
